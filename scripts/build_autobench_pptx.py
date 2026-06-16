"""Generate a premium, Mastercard-branded PPTX of the Autobench v3.0 deck.

Content is sourced from ``docs/Autobench_v3.0_Precision_Compliance.pdf`` and
restyled with the visual identity of the Mastercard "Acquiring Performance
Map / New Data Strategy" template
(``docs/Acquiring_Performance_Map_New_Data_Strategy-v2 1.pdf``).

Design notes
------------
* Fully native PowerPoint objects (shapes, text, tables, freeform charts) — no
  rasterised slide images, so every element stays editable.
* Typography is set in Inter Display / Inter / JetBrains Mono with Font Awesome
  for iconography; all of these are **embedded** into the file (see
  ``pptx_embed_fonts``) so the deck renders identically without local installs.
* A restrained palette (warm neutrals + a single Mastercard-orange accent) and
  a strict layout grid replace the heavier first draft.

Usage::

    py scripts/build_autobench_pptx.py [output.pptx]
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

sys.path.insert(0, str(Path(__file__).resolve().parent))
from pptx_embed_fonts import embed_fonts  # noqa: E402

# --------------------------------------------------------------------------- #
# Typefaces (all embedded)
# --------------------------------------------------------------------------- #
HERO = "Inter Display Light"
DISPLAY = "Inter Display"          # regular slot = SemiBold, bold slot = Bold
DISPLAY_MED = "Inter Display Medium"
SANS = "Inter"                     # regular + bold
SANS_MED = "Inter Medium"
SANS_SB = "Inter SemiBold"
MONO = "JetBrains Mono"
ICON = "Font Awesome 6 Free Solid"

FONT_DIR = Path(__file__).resolve().parent.parent / "assets" / "fonts"
EMBED_MAP = {
    HERO: {"regular": "InterDisplay-Light.ttf"},
    DISPLAY: {"regular": "InterDisplay-SemiBold.ttf",
              "bold": "InterDisplay-Bold.ttf"},
    DISPLAY_MED: {"regular": "InterDisplay-Medium.ttf"},
    SANS: {"regular": "Inter-Regular.ttf", "bold": "Inter-Bold.ttf"},
    SANS_MED: {"regular": "Inter-Medium.ttf"},
    SANS_SB: {"regular": "Inter-SemiBold.ttf"},
    MONO: {"regular": "JetBrainsMono-Regular.ttf",
           "bold": "JetBrainsMono-Bold.ttf"},
    ICON: {"regular": "FontAwesome6Free-Solid.ttf"},
}

# Font Awesome 6 (Solid) codepoints
FA = {
    "terminal": "\uf120", "desktop": "\uf108", "gears": "\uf085",
    "sliders": "\uf1de", "shield": "\uf3ed", "gauge": "\uf625",
    "excel": "\uf1c3", "csv": "\uf6dd", "archive": "\uf187",
    "database": "\uf1c0", "chip": "\uf2db", "filter": "\uf0b0",
    "chart": "\uf201", "scale": "\uf24e", "compass": "\uf568",
    "magnifier": "\uf002", "check": "\uf058", "warn": "\uf071",
    "ban": "\uf05e", "bolt": "\uf0e7", "branch": "\uf126",
    "clipboard": "\uf46c", "lock": "\uf023", "layers": "\uf5fd",
    "wand": "\uf72b", "table": "\uf0ce",
}

# --------------------------------------------------------------------------- #
# Palette (sampled from / harmonised with the Mastercard template)
# --------------------------------------------------------------------------- #
INK = RGBColor(0x1A, 0x1A, 0x1A)
INK2 = RGBColor(0x3C, 0x3F, 0x44)
GRAY = RGBColor(0x65, 0x6A, 0x71)
MUTE = RGBColor(0x9C, 0xA0, 0xA6)
HAIR = RGBColor(0xE6, 0xE2, 0xDC)
HAIR2 = RGBColor(0xEF, 0xEC, 0xE7)
PANEL = RGBColor(0xFB, 0xFA, 0xF8)
PANEL2 = RGBColor(0xF3, 0xF1, 0xEC)
ORANGE = RGBColor(0xFF, 0x67, 0x1B)
ORANGE_DK = RGBColor(0xD8, 0x4E, 0x0C)
ORANGE_SOFT = RGBColor(0xFF, 0xED, 0xE2)
RED = RGBColor(0xEB, 0x00, 0x1B)
AMBER = RGBColor(0xF7, 0x9E, 0x1B)
OVERLAP = RGBColor(0xFF, 0x5F, 0x00)
DEEP = RGBColor(0x1D, 0x27, 0x35)
DEEP_SOFT = RGBColor(0xEC, 0xEE, 0xF2)
CREAM = RGBColor(0xF4, 0xF2, 0xEF)
CREAM2 = RGBColor(0xFB, 0xFA, 0xF7)
DARK = RGBColor(0x15, 0x17, 0x1C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# --------------------------------------------------------------------------- #
# Geometry — 16:9 widescreen on a tidy grid
# --------------------------------------------------------------------------- #
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)
MARGIN = Pt(48)
CONTENT_W = SLIDE_W - MARGIN * 2
CONFIDENTIAL = "\u00a9 2025 Mastercard. Proprietary and Confidential."
TOTAL = 13


# --------------------------------------------------------------------------- #
# Low-level primitives
# --------------------------------------------------------------------------- #
def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _noline(shape):
    shape.line.fill.background()


def _stroke(shape, color, w=Pt(1)):
    shape.line.color.rgb = color
    shape.line.width = w


def _flat(shape):
    """Remove the theme style ref (drop shadow / preset fill)."""
    st = shape._element.find(qn("p:style"))
    if st is not None:
        shape._element.remove(st)
    sp = shape._element.spPr
    if sp.find(qn("a:effectLst")) is None:
        sp.append(sp.makeelement(qn("a:effectLst"), {}))


def _alpha(shape, pct):
    """pct = opacity percent (100 = opaque)."""
    srgb = shape.fill.fore_color._xFill.find(qn("a:srgbClr"))
    if srgb is not None:
        srgb.append(srgb.makeelement(qn("a:alpha"), {"val": str(int(pct * 1000))}))


def _shadow(shape, blur=Pt(9), dist=Pt(3), alpha=22, direction=5400000):
    """Soft, subtle drop shadow for elevated surfaces."""
    sp = shape._element.spPr
    old = sp.find(qn("a:effectLst"))
    if old is not None:
        sp.remove(old)
    eff = sp.makeelement(qn("a:effectLst"), {})
    sh = eff.makeelement(qn("a:outerShdw"), {
        "blurRad": str(int(blur)), "dist": str(int(dist)),
        "dir": str(direction), "rotWithShape": "0",
    })
    clr = sh.makeelement(qn("a:srgbClr"), {"val": "1A1A1A"})
    clr.append(clr.makeelement(qn("a:alpha"), {"val": str(alpha * 1000)}))
    sh.append(clr)
    eff.append(sh)
    sp.append(eff)


def rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(1), radius=None,
         shadow=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE,
        int(x), int(y), int(w), int(h))
    if fill is None:
        shp.fill.background()
    else:
        _fill(shp, fill)
    if line is None:
        _noline(shp)
    else:
        _stroke(shp, line, line_w)
    if shadow:
        _shadow(shp)
        st = shp._element.find(qn("p:style"))
        if st is not None:
            shp._element.remove(st)
    else:
        _flat(shp)
    if radius is not None:
        try:
            shp.adjustments[0] = radius
        except (IndexError, KeyError):
            pass
    return shp


def oval(slide, x, y, w, h, fill, line=None, line_w=Pt(1)):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(x), int(y), int(w), int(h))
    _fill(shp, fill)
    if line is None:
        _noline(shp)
    else:
        _stroke(shp, line, line_w)
    _flat(shp)
    return shp


def line_seg(slide, x1, y1, x2, y2, color, w=Pt(1), dash=None):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, int(x1), int(y1),
                                   int(x2), int(y2))
    c.line.color.rgb = color
    c.line.width = w
    if dash:
        ln = c.line._get_or_add_ln()
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": dash}))
    st = c._element.find(qn("p:style"))
    if st is not None:
        c._element.remove(st)
    return c


def free_poly(slide, pts, fill=None, line=None, line_w=Pt(1)):
    fb = slide.shapes.build_freeform(float(pts[0][0]), float(pts[0][1]),
                                     scale=1.0)
    fb.add_line_segments([(float(x), float(y)) for x, y in pts[1:]],
                         close=True)
    shp = fb.convert_to_shape()
    if fill is None:
        shp.fill.background()
    else:
        _fill(shp, fill)
    if line is None:
        _noline(shp)
    else:
        _stroke(shp, line, line_w)
    _flat(shp)
    return shp


# --------------------------------------------------------------------------- #
# Text
# --------------------------------------------------------------------------- #
def _apply(p, runs, size, color, font, bold, lh, sa, sb, align, track):
    if align is not None:
        p.alignment = align
    if lh is not None:
        p.line_spacing = lh
    if sa is not None:
        p.space_after = sa
    if sb is not None:
        p.space_before = sb
    for text, o in runs:
        r = p.add_run()
        r.text = text
        f = r.font
        f.size = o.get("size", size)
        f.name = o.get("font", font)
        f.bold = o.get("bold", bold)
        f.italic = o.get("italic", False)
        f.color.rgb = o.get("color", color)
        tr = o.get("track", track)
        if tr:
            r._r.get_or_add_rPr().set("spc", str(int(tr * 100)))


def text(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ("left", "right", "top", "bottom"):
        setattr(tf, f"margin_{m}", 0)
    for i, pa in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _apply(p, pa["runs"], pa.get("size", Pt(13)), pa.get("color", INK),
               pa.get("font", SANS), pa.get("bold", False),
               pa.get("lh"), pa.get("sa"), pa.get("sb"),
               pa.get("align"), pa.get("track"))
    return tb


def run(t, **o):
    return [(t, o)]


def para(runs, **kw):
    kw["runs"] = runs
    return kw


# --------------------------------------------------------------------------- #
# Brand furniture
# --------------------------------------------------------------------------- #
def mastercard_mark(slide, right, top, d):
    """Accurate two-circle Mastercard symbol with the interlock lens."""
    sep = d * 0.52
    total = sep + d
    rx = right - total
    ax = rx + sep
    oval(slide, rx, top, d, d, RED)
    oval(slide, ax, top, d, d, AMBER)
    lens_w = d - sep
    lens_h = d * 0.855
    oval(slide, ax, top + (d - lens_h) / 2, lens_w, lens_h, OVERLAP)


def icon_token(slide, x, y, size, glyph, fg=ORANGE, bg=ORANGE_SOFT,
               circle=False, gsize=None):
    if bg is not None:
        if circle:
            oval(slide, x, y, size, size, bg)
        else:
            rect(slide, x, y, size, size, fill=bg, radius=0.26)
    text(slide, x, y, size, size,
         [para(run(glyph, font=ICON, color=fg, size=gsize or Pt(size / 12700 * 0.52)),
               align=PP_ALIGN.CENTER)],
         anchor=MSO_ANCHOR.MIDDLE)


def eyebrow(slide, x, y, label, color=ORANGE, w=None):
    text(slide, x, y, w or Pt(500), Pt(16),
         [para(run(label.upper(), font=SANS_SB, color=color, size=Pt(10.5),
                   track=1.6))])


def header(slide, eb, headline, dek=None, hl_size=Pt(23)):
    eyebrow(slide, MARGIN, Pt(40), eb)
    text(slide, MARGIN, Pt(58), CONTENT_W, Pt(74),
         [para(run(headline, font=DISPLAY, color=INK, size=hl_size), lh=1.06)])
    rule_y = Pt(128) if dek is None else Pt(134)
    if dek:
        text(slide, MARGIN, Pt(102), CONTENT_W, Pt(26),
             [para(run(dek, font=SANS, color=GRAY, size=Pt(13.5)), lh=1.12)])
    line_seg(slide, MARGIN, rule_y, SLIDE_W - MARGIN, rule_y, HAIR, Pt(1))
    rect(slide, MARGIN, rule_y - Pt(1.2), Pt(54), Pt(2.4), fill=ORANGE)
    return rule_y + Pt(24)


def footer(slide, page, dark=False):
    col = MUTE if not dark else RGBColor(0x8A, 0x8E, 0x96)
    text(slide, MARGIN, SLIDE_H - Pt(30), Pt(200), Pt(16),
         [para([(f"{page:02d}", {"font": SANS_SB, "color": ORANGE if not dark else AMBER}),
                (f"  /  {TOTAL:02d}", {"font": SANS_MED, "color": col})],
               size=Pt(9))],
         anchor=MSO_ANCHOR.MIDDLE)
    box = slide.shapes.add_textbox(SLIDE_W - Pt(22), Pt(130), Pt(14),
                                   SLIDE_H - Pt(280))
    tf = box.text_frame
    tf.word_wrap = False
    for m in ("left", "right", "top", "bottom"):
        setattr(tf, f"margin_{m}", 0)
    p = tf.paragraphs[0]
    rr = p.add_run()
    rr.text = CONFIDENTIAL
    rr.font.size = Pt(6.5)
    rr.font.name = SANS_MED
    rr.font.color.rgb = col
    tf._txBody.find(qn("a:bodyPr")).set("vert", "vert270")
    mastercard_mark(slide, SLIDE_W - MARGIN, SLIDE_H - Pt(48), Pt(26))


def bg(slide, color=WHITE):
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=color)


def bullets(items, size=Pt(12.5), color=INK2, lh=1.18, sa=Pt(7),
            marker=ORANGE):
    out = []
    for runs in items:
        out.append(para([("\u2014  ", {"font": SANS_SB, "color": marker})]
                        + list(runs), size=size, color=color, font=SANS,
                        lh=lh, sa=sa))
    return out


# --------------------------------------------------------------------------- #
# Reusable composite components
# --------------------------------------------------------------------------- #
def feature_card(slide, x, y, w, h, glyph, title, body_runs, accent=ORANGE,
                 shadow=True, title_size=Pt(15)):
    rect(slide, x, y, w, h, fill=WHITE, line=HAIR, line_w=Pt(1), radius=0.05,
         shadow=shadow)
    pad = Pt(20)
    icon_token(slide, x + pad, y + pad, Pt(38), glyph, fg=accent,
               bg=ORANGE_SOFT if accent == ORANGE else DEEP_SOFT)
    text(slide, x + pad, y + pad + Pt(50), w - pad * 2, Pt(46),
         [para(run(title, font=DISPLAY, color=INK, size=title_size), lh=1.02)])
    text(slide, x + pad, y + pad + Pt(96), w - pad * 2, h - pad * 2 - Pt(96),
         body_runs)


def chip(slide, x, y, label, fg, border, solid=False, mono=True):
    w = Pt(18 + len(label) * 5.6)
    rect(slide, x, y, w, Pt(24), fill=fg if solid else WHITE,
         line=border, line_w=Pt(1.1), radius=0.5)
    text(slide, x, y, w, Pt(24),
         [para(run(label, font=MONO if mono else SANS_SB,
                   color=WHITE if solid else fg, size=Pt(9)),
               align=PP_ALIGN.CENTER)],
         anchor=MSO_ANCHOR.MIDDLE)
    return w


# --------------------------------------------------------------------------- #
# Slides
# --------------------------------------------------------------------------- #
def s_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def s_title(prs):
    s = s_blank(prs)
    bg(s, CREAM)
    # signature arc field
    big = oval(s, Emu(int(SLIDE_W * 0.34)), Emu(int(-SLIDE_H * 0.58)),
               Emu(int(SLIDE_W * 1.0)), Emu(int(SLIDE_H * 2.0)), CREAM2)
    _flat(big)
    ring = s.shapes.add_shape(MSO_SHAPE.OVAL, Emu(int(SLIDE_W * 0.39)),
                              Emu(int(-SLIDE_H * 0.5)), Emu(int(SLIDE_W * 0.92)),
                              Emu(int(SLIDE_H * 1.84)))
    ring.fill.background()
    _stroke(ring, HAIR, Pt(1))
    _flat(ring)
    mastercard_mark(s, SLIDE_W - MARGIN, Pt(40), Pt(40))
    text(s, MARGIN, Pt(150), Pt(560), Pt(20),
         [para(run("MASTERCARD ADVISORS  \u2022  PERFORMANCE ANALYTICS",
                   font=SANS_SB, color=ORANGE, size=Pt(11), track=2.0))])
    text(s, MARGIN - Pt(2), Pt(212), Pt(820), Pt(180),
         [para(run("Autobench v3.0", font=HERO, color=INK, size=Pt(58)),
               lh=1.0),
          para(run("The Precision Compliance Engine", font=DISPLAY, color=INK,
                   size=Pt(34)), lh=1.05, sb=Pt(6))])
    line_seg(s, MARGIN, Pt(404), MARGIN + Pt(360), Pt(404), ORANGE, Pt(2.2))
    text(s, MARGIN, Pt(420), Pt(620), Pt(60),
         [para(run("Next-generation benchmarking. Minimal distortion. "
                   "Absolute control.", font=DISPLAY_MED, color=INK2,
                   size=Pt(16)), lh=1.2)])
    text(s, MARGIN, Pt(486), Pt(620), Pt(24),
         [para(run("Executive Overview & Operational Capabilities",
                   font=SANS, color=GRAY, size=Pt(12)))])
    return s


def s_paradox(prs, page):
    s = s_blank(prs)
    bg(s)
    top = header(
        s, "The compliance paradox",
        "Compliant benchmarking forces a trade-off between privacy and "
        "market truth.",
        "Autobench v3.0 resolves it mathematically \u2014 not manually.")
    gap = Pt(120)
    cw = Emu(int((CONTENT_W - gap) / 2))
    rx = MARGIN + cw + gap
    cards_h = Pt(212)

    def panel(x, glyph, kicker, kcolor, title, items):
        rect(s, x, top, cw, cards_h, fill=PANEL, line=HAIR, line_w=Pt(1),
             radius=0.045)
        icon_token(s, x + Pt(22), top + Pt(22), Pt(38), glyph, fg=kcolor,
                   bg=ORANGE_SOFT if kcolor == ORANGE else DEEP_SOFT)
        eyebrow(s, x + Pt(72), top + Pt(26), kicker, color=kcolor)
        text(s, x + Pt(72), top + Pt(42), cw - Pt(94), Pt(24),
             [para(run(title, font=DISPLAY, color=INK, size=Pt(15.5)))])
        text(s, x + Pt(24), top + Pt(86), cw - Pt(48), cards_h - Pt(100),
             bullets([run(i, color=INK2) for i in items], size=Pt(12.5),
                     sa=Pt(7), marker=kcolor))

    panel(MARGIN, FA["lock"], "The risk", DEEP, "Control 3 constraints", [
        "Strict concentration caps (no peer above 25% share).",
        "Minimum participant thresholds per category.",
        "Zero-tolerance regulatory environment.",
        "Manual weighting silently destroys data fidelity.",
    ])
    panel(rx, FA["compass"], "The goal", ORANGE, "Market reality", [
        "Maintain exact market positioning.",
        "Preserve rank integrity across every dimension.",
        "Deliver actionable merchant and issuer performance.",
    ])
    # central imperative arrow
    ay = top + cards_h / 2 - Pt(20)
    arr = s.shapes.add_shape(MSO_SHAPE.CHEVRON, MARGIN + cw + Pt(28), ay,
                             gap - Pt(56), Pt(40))
    _fill(arr, ORANGE)
    _noline(arr)
    _flat(arr)
    text(s, MARGIN + cw, top + cards_h / 2 + Pt(26), gap, Pt(18),
         [para(run("AUTOBENCH", font=SANS_SB, color=ORANGE_DK, size=Pt(8),
                   track=1.2), align=PP_ALIGN.CENTER)])

    band_y = top + cards_h + Pt(34)
    rect(s, MARGIN, band_y, CONTENT_W, Pt(74), fill=ORANGE, radius=0.08)
    icon_token(s, MARGIN + Pt(26), band_y + Pt(20), Pt(34), FA["wand"],
               fg=WHITE, bg=None)
    text(s, MARGIN + Pt(74), band_y, CONTENT_W - Pt(120), Pt(74),
         [para([("The solution.  ", {"font": DISPLAY, "color": WHITE,
                                     "size": Pt(14)}),
                ("Advanced Linear Programming enforces absolute privacy "
                 "compliance while mathematically minimizing distortion.",
                 {"font": SANS, "color": WHITE, "size": Pt(13.5)})], lh=1.14)],
         anchor=MSO_ANCHOR.MIDDLE)
    footer(s, page)
    return s


def s_architecture(prs, page):
    s = s_blank(prs)
    bg(s)
    top = header(
        s, "Architecture",
        "Three coordinated phases turn raw data into audit-ready "
        "intelligence.")
    h = Pt(296)
    aw = Pt(40)
    cw = Emu(int((CONTENT_W - aw * 2) / 3))
    x0 = MARGIN
    x1 = x0 + cw + aw
    x2 = x1 + cw + aw

    def phase(x, idx, glyph, name, sub, rows, primary=False):
        rect(s, x, top, cw, h, fill=PANEL if not primary else WHITE,
             line=ORANGE if primary else HAIR,
             line_w=Pt(1.5) if primary else Pt(1), radius=0.04,
             shadow=primary)
        rect(s, x, top, cw, Pt(4), fill=ORANGE if primary else DEEP,
             radius=None)
        icon_token(s, x + Pt(20), top + Pt(22), Pt(40), glyph,
                   fg=ORANGE if primary else DEEP,
                   bg=ORANGE_SOFT if primary else DEEP_SOFT)
        eyebrow(s, x + Pt(72), top + Pt(24), f"Phase {idx}",
                color=ORANGE if primary else DEEP)
        text(s, x + Pt(72), top + Pt(40), cw - Pt(90), Pt(40),
             [para(run(name, font=DISPLAY, color=INK, size=Pt(15)), lh=1.0)])
        if sub:
            text(s, x + Pt(20), top + Pt(82), cw - Pt(40), Pt(18),
                 [para(run(sub, font=MONO, color=GRAY, size=Pt(9)))])
        ry = top + Pt(108)
        iw = cw - Pt(40)
        rh = Pt(48)
        for label in rows:
            rect(s, x + Pt(20), ry, iw, rh, fill=WHITE if not primary else PANEL,
                 line=HAIR2, line_w=Pt(1), radius=0.12)
            text(s, x + Pt(32), ry, iw - Pt(24), rh,
                 [para(run(label, font=SANS_MED, color=INK2, size=Pt(11)),
                       lh=1.02)],
                 anchor=MSO_ANCHOR.MIDDLE)
            ry += rh + Pt(10)

    phase(x0, 1, FA["terminal"], "Ingestion Interfaces", None, [
        "CLI \u00b7 automated pipelines",
        "TUI \u00b7 analyst workbench",
    ])
    phase(x1, 2, FA["gears"], "Core Engine", "core/dimensional_analyzer.py", [
        "Normalization & schema validation",
        "Control 3.2 privacy validator",
        "Global LP optimization \u00b7 HiGHS",
    ], primary=True)
    phase(x2, 3, FA["archive"], "Output Artifacts", None, [
        "Excel reporting & diagnostics",
        "Balanced CSV \u00b7 BI-ready",
        "Comprehensive audit packages",
    ])
    for ax in (x0 + cw, x1 + cw):
        oval(s, ax + (aw - Pt(26)) / 2, top + h / 2 - Pt(13), Pt(26), Pt(26),
             WHITE, line=HAIR, line_w=Pt(1))
        text(s, ax, top + h / 2 - Pt(13), aw, Pt(26),
             [para(run("\uf061", font=ICON, color=ORANGE, size=Pt(10)),
                   align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    footer(s, page)
    return s


def s_solver(prs, page):
    s = s_blank(prs)
    bg(s)
    top = header(
        s, "Algorithmic resilience",
        "A three-tier solver cascade guarantees an answer \u2014 even when "
        "strict optimization fails.")
    levels = [
        (FA["gauge"], "01", "Global LP Optimization", ORANGE, [
            "Solves all dimensions simultaneously.",
            "Minimizes deviation from 1.0 + rank penalty + slack.",
            "Zero privacy violations, lowest system-wide distortion.",
        ], "If structurally infeasible"),
        (FA["filter"], "02", "Subset Search & Per-Dimension LP", DEEP, [
            "Greedily isolates the largest feasible subset of dimensions.",
            "Routes conflicting dimensions to independent solvers.",
        ], "If strict LP fails"),
        (FA["wand"], "03", "Bayesian / Heuristic Fallback", GRAY, [
            "Engages L-BFGS-B (scipy.minimize) on edge-case categories.",
            "Guarantees a usable output under extreme data sparsity.",
        ], None),
    ]
    x = MARGIN
    w = CONTENT_W
    y = top
    hs = [Pt(94), Pt(80), Pt(94)]
    for (glyph, num, title, accent, items, cond), hh in zip(levels, hs):
        rect(s, x, y, w, hh, fill=PANEL, line=HAIR, line_w=Pt(1), radius=0.06)
        rect(s, x, y, Pt(5), hh, fill=accent)
        icon_token(s, x + Pt(24), y + (hh - Pt(40)) / 2, Pt(40), glyph,
                   fg=accent, bg=ORANGE_SOFT if accent == ORANGE else DEEP_SOFT)
        text(s, x + Pt(80), y + Pt(16), Pt(60), Pt(40),
             [para(run(num, font=HERO, color=accent, size=Pt(26)))])
        text(s, x + Pt(150), y + Pt(15), w - Pt(180), Pt(26),
             [para(run(title, font=DISPLAY, color=INK, size=Pt(15.5)))])
        text(s, x + Pt(150), y + Pt(42), w - Pt(190), hh - Pt(50),
             bullets([run(i, color=INK2) for i in items], size=Pt(11.5),
                     sa=Pt(3), lh=1.12, marker=accent))
        if cond:
            cy = y + hh + Pt(2)
            text(s, x + Pt(150), cy, Pt(420), Pt(20),
                 [para([("\uf063  ", {"font": ICON, "color": ORANGE,
                                      "size": Pt(8)}),
                        (cond.upper(), {"font": SANS_SB, "color": ORANGE_DK,
                                        "size": Pt(8.5), "track": 1.0})])])
        y += hh + Pt(24)
    footer(s, page)
    return s


def styled_table_caps(s, x, y, w):
    headers = ["Rule", "Min peers", "Max concentration", "Trigger condition"]
    rows = [
        ("5 / 25", "5", "25%", "Baseline standard", False),
        ("6 / 30", "6", "30%", "\u2265 3 participants at \u2265 7%", False),
        ("7 / 35", "7", "35%", "\u2265 2 at \u2265 15%, plus \u2265 1 at \u2265 8%", False),
        ("10 / 40", "10", "40%", "\u2265 2 at \u2265 20%, plus \u2265 1 at \u2265 10%", False),
        ("4 / 35", "4", "35%", "Merchant benchmarking mode only", True),
    ]
    widths = [Pt(150), Pt(130), Pt(220), Pt(396)]
    head_h = Pt(40)
    row_h = Pt(50)
    # header
    text_x = x
    for htext, cwid in zip(headers, widths):
        text(s, text_x + Pt(2), y, cwid - Pt(4), head_h,
             [para(run(htext.upper(), font=SANS_SB, color=GRAY, size=Pt(10),
                       track=0.8))], anchor=MSO_ANCHOR.MIDDLE)
        text_x += cwid
    line_seg(s, x, y + head_h, x + w, y + head_h, INK, Pt(1.6))
    ry = y + head_h
    for r in rows:
        merchant = r[4]
        if merchant:
            rect(s, x, ry, w, row_h, fill=ORANGE_SOFT)
        cxx = x
        for ci, (val, cwid) in enumerate(zip(r[:4], widths)):
            if ci == 0:
                rect(s, cxx + Pt(2), ry + (row_h - Pt(26)) / 2,
                     Pt(64), Pt(26),
                     fill=ORANGE if not merchant else WHITE,
                     line=ORANGE if merchant else None,
                     line_w=Pt(1.2), radius=0.5)
                text(s, cxx + Pt(2), ry + (row_h - Pt(26)) / 2, Pt(64), Pt(26),
                     [para(run(val, font=SANS_SB,
                               color=WHITE if not merchant else ORANGE_DK,
                               size=Pt(10.5)), align=PP_ALIGN.CENTER)],
                     anchor=MSO_ANCHOR.MIDDLE)
            else:
                fnt = MONO if ci in (1, 2) else SANS
                col = INK if ci != 3 else INK2
                text(s, cxx + Pt(4), ry, cwid - Pt(8), row_h,
                     [para(run(val, font=fnt, color=col,
                               size=Pt(12) if ci != 3 else Pt(12.5)))],
                     anchor=MSO_ANCHOR.MIDDLE)
            cxx += cwid
        line_seg(s, x, ry + row_h, x + w, ry + row_h, HAIR2, Pt(1))
        ry += row_h


def s_caps(prs, page):
    s = s_blank(prs)
    bg(s)
    top = header(
        s, "Control 3.2",
        "Privacy caps are selected automatically from the available peer "
        "data.",
        "No manual configuration \u2014 the engine applies the correct "
        "regulatory rule.")
    styled_table_caps(s, MARGIN, top + Pt(6), CONTENT_W)
    footer(s, page)
    return s


def s_presets(prs, page):
    s = s_blank(prs)
    bg(s)
    top = header(
        s, "Strategic presets",
        "Out-of-the-box configurations mapped to specific business intents.")
    presets = [
        (FA["scale"], "Balanced Default", "balanced_default.yaml",
         "Day-to-day general business analysis.",
         "Optimal mix of speed, compliance, and accuracy."),
        (FA["shield"], "Compliance Strict", "compliance_strict.yaml",
         "Regulatory and audit-first submissions.",
         "Zero tolerance for threshold deviations; robust audit logs."),
        (FA["chart"], "Strategic Consistency", "strategic_consistency.yaml",
         "Executive dashboards and KPI tracking.",
         "Consistent global weighting behavior over time."),
        (FA["magnifier"], "Research Exploratory", "research_exploratory.yaml",
         "Hard, sparse, or highly skewed datasets.",
         "Higher flexibility to surface underlying market trends."),
    ]
    h = Pt(290)
    gap = Pt(22)
    n = len(presets)
    cw = Emu(int((CONTENT_W - gap * (n - 1)) / n))
    x = MARGIN
    for glyph, name, yaml, use, posture in presets:
        rect(s, x, top, cw, h, fill=WHITE, line=HAIR, line_w=Pt(1),
             radius=0.05, shadow=True)
        rect(s, x, top, cw, Pt(4), fill=ORANGE)
        pad = Pt(18)
        icon_token(s, x + pad, top + Pt(24), Pt(42), glyph, fg=ORANGE)
        text(s, x + pad, top + Pt(80), cw - pad * 2, Pt(44),
             [para(run(name, font=DISPLAY, color=INK, size=Pt(15)), lh=1.0)])
        text(s, x + pad, top + Pt(124), cw - pad * 2, Pt(18),
             [para(run(yaml, font=MONO, color=ORANGE_DK, size=Pt(9.5)))])
        line_seg(s, x + pad, top + Pt(150), x + cw - pad, top + Pt(150),
                 HAIR2, Pt(1))
        text(s, x + pad, top + Pt(162), cw - pad * 2, h - Pt(176),
             [para([("Use case\n", {"font": SANS_SB, "color": GRAY,
                                    "size": Pt(9.5), "track": 0.6}),
                    (use, {"font": SANS, "color": INK2, "size": Pt(11.5)})],
                   lh=1.12, sa=Pt(12)),
              para([("Posture\n", {"font": SANS_SB, "color": GRAY,
                                   "size": Pt(9.5), "track": 0.6}),
                    (posture, {"font": SANS, "color": INK2, "size": Pt(11.5)})],
                   lh=1.12)])
        x += cw + gap
    footer(s, page)
    return s


def s_tui(prs, page):
    s = s_blank(prs)
    bg(s)
    top = header(
        s, "Analyst experience",
        "A robust TUI eliminates pure-CLI friction.",
        "Validation-first, with live diagnostics and instant overrides.")
    win_w = Pt(430)
    win_h = Pt(300)
    win_x = MARGIN + (CONTENT_W - win_w) / 2
    win_y = top + Pt(6)
    rect(s, win_x, win_y, win_w, win_h, fill=WHITE, line=HAIR, line_w=Pt(1),
         radius=0.03, shadow=True)
    rect(s, win_x, win_y, win_w, Pt(30), fill=DEEP, radius=0.03)
    rect(s, win_x, win_y + Pt(15), win_w, Pt(15), fill=DEEP)
    for i, c in enumerate((RED, AMBER, RGBColor(0x3F, 0xB9, 0x50))):
        oval(s, win_x + Pt(12) + i * Pt(14), win_y + Pt(11), Pt(8), Pt(8), c)
    text(s, win_x, win_y, win_w, Pt(30),
         [para(run("analyst validation dashboard \u00b7 v3.0", font=MONO,
                   color=WHITE, size=Pt(9)), align=PP_ALIGN.CENTER)],
         anchor=MSO_ANCHOR.MIDDLE)
    # status strip
    sy = win_y + Pt(40)
    rect(s, win_x + Pt(12), sy, win_w - Pt(24), Pt(24), fill=PANEL2,
         radius=0.2)
    for i, (lab, st, col) in enumerate([
            ("SCHEMA", "PASS", RGBColor(0x2E, 0x7D, 0x32)),
            ("NULLS", "PASS", RGBColor(0x2E, 0x7D, 0x32)),
            ("PEERS", "VALIDATING", ORANGE_DK)]):
        bx = win_x + Pt(20) + i * Pt((430 - 40) / 3)
        text(s, bx, sy, Pt(130), Pt(24),
             [para([(lab + " ", {"font": MONO, "color": GRAY, "size": Pt(8)}),
                    (st, {"font": MONO, "color": col, "size": Pt(8),
                          "bold": True})])], anchor=MSO_ANCHOR.MIDDLE)
    # two panes
    py = sy + Pt(32)
    pane_h = Pt(120)
    rect(s, win_x + Pt(12), py, Pt(150), pane_h, fill=PANEL, line=HAIR2,
         line_w=Pt(1), radius=0.04)
    text(s, win_x + Pt(22), py + Pt(10), Pt(132), pane_h - Pt(16),
         [para(run("CONFIGURATION", font=MONO, color=MUTE, size=Pt(7.5),
                   track=0.5), sa=Pt(8)),
          para(run("entity   \u25be", font=MONO, color=INK2, size=Pt(8.5)),
               sa=Pt(6)),
          para(run("metric   \u25be", font=MONO, color=INK2, size=Pt(8.5)),
               sa=Pt(6)),
          para(run("dimension\u25be", font=MONO, color=INK2, size=Pt(8.5)))])
    dx = win_x + Pt(170)
    dw = win_w - Pt(182)
    rect(s, dx, py, dw, pane_h, fill=PANEL, line=HAIR2, line_w=Pt(1),
         radius=0.04)
    rows = [("E123", "2024-05-20", "150.5", "VALID"),
            ("E124", "2024-05-20", "150.5", "VALID"),
            ("E456", "2024-05-20", "  \u2014", "ERROR")]
    text(s, dx + Pt(10), py + Pt(10), dw - Pt(20), Pt(16),
         [para(run("DATA PREVIEW", font=MONO, color=MUTE, size=Pt(7.5),
                   track=0.5))])
    for i, (a, b, c, d) in enumerate(rows):
        col = RGBColor(0xB0, 0x2A, 0x1B) if d == "ERROR" else INK2
        text(s, dx + Pt(10), py + Pt(30) + i * Pt(20), dw - Pt(20), Pt(18),
             [para([(f"{a}   {b}   {c}   ", {"font": MONO, "color": INK2,
                                            "size": Pt(8.5)}),
                    (d, {"font": MONO, "color": col, "size": Pt(8.5),
                         "bold": True})])])
    # log
    ly = py + pane_h + Pt(10)
    rect(s, win_x + Pt(12), ly, win_w - Pt(24), Pt(40), fill=DARK, radius=0.05)
    text(s, win_x + Pt(22), ly, win_w - Pt(40), Pt(40),
         [para([("WARN ", {"font": MONO, "color": AMBER, "size": Pt(7.5),
                           "bold": True}),
                ("schema mismatch \u2192 auto-fixing\u2026",
                 {"font": MONO, "color": RGBColor(0xC8, 0xCB, 0xD0),
                  "size": Pt(7.5)})], sa=Pt(2)),
          para([("ERR  ", {"font": MONO, "color": RGBColor(0xFF, 0x6B, 0x5B),
                           "size": Pt(7.5), "bold": True}),
                ("peer constraint E456 < min_peers",
                 {"font": MONO, "color": RGBColor(0xC8, 0xCB, 0xD0),
                  "size": Pt(7.5)})])],
         anchor=MSO_ANCHOR.MIDDLE)

    feats_l = [(FA["shield"], "Validation-first", "Catches nulls, schema "
                "errors, and peer-constraint breaches before any heavy run."),
               (FA["sliders"], "Intuitive config", "Dropdown selection of "
                "entities, metrics, and dimensions.")]
    feats_r = [(FA["chart"], "Real-time diagnostics", "Integrated log handling "
                "keeps the UI responsive during background work."),
               (FA["bolt"], "Instant overrides", "Toggle advanced panels "
                "(Ctrl+A) to inject custom rules on the fly.")]

    def feat(x, y, glyph, t, b, w):
        icon_token(s, x, y, Pt(30), glyph, fg=ORANGE)
        text(s, x + Pt(40), y - Pt(2), w - Pt(40), Pt(20),
             [para(run(t, font=DISPLAY, color=INK, size=Pt(12.5)))])
        text(s, x + Pt(40), y + Pt(20), w - Pt(40), Pt(70),
             [para(run(b, font=SANS, color=GRAY, size=Pt(10.5)), lh=1.14)])

    lw = win_x - MARGIN - Pt(16)
    feat(MARGIN, win_y + Pt(18), *feats_l[0], lw)
    feat(MARGIN, win_y + Pt(150), *feats_l[1], lw)
    rxx = win_x + win_w + Pt(16)
    rw = SLIDE_W - MARGIN - rxx
    feat(rxx, win_y + Pt(18), *feats_r[0], rw)
    feat(rxx, win_y + Pt(150), *feats_r[1], rw)
    footer(s, page)
    return s


def s_resource(prs, page):
    s = s_blank(prs)
    bg(s)
    top = header(
        s, "Resource management",
        "Adaptive batching and lean mode scale to massive, memory-limited "
        "environments.")
    steps = [
        (FA["magnifier"], "Estimate & Project",
         "Reads headers, drops unused columns, projects memory load."),
        (FA["layers"], "Adaptive Chunking",
         "Streams heavy CSVs in 100k-row chunks."),
        (FA["filter"], "Pre-Aggregation",
         "Consolidates duplicate entity / dimension / time rows in flight."),
        (FA["chip"], "Lean Execution",
         "Disables secondary artifacts to focus 100% on global optimization."),
    ]
    h = Pt(220)
    gap = Pt(22)
    n = len(steps)
    cw = Emu(int((CONTENT_W - gap * (n - 1)) / n))
    x = MARGIN
    for i, (glyph, name, desc) in enumerate(steps):
        rect(s, x, top, cw, h, fill=PANEL, line=HAIR, line_w=Pt(1),
             radius=0.05)
        pad = Pt(18)
        icon_token(s, x + pad, top + Pt(20), Pt(40), glyph, fg=ORANGE)
        text(s, x + cw - Pt(54), top + Pt(16), Pt(40), Pt(40),
             [para(run(f"{i + 1}", font=HERO, color=HAIR, size=Pt(34)),
                   align=PP_ALIGN.RIGHT)])
        text(s, x + pad, top + Pt(74), cw - pad * 2, Pt(40),
             [para(run(name, font=DISPLAY, color=INK, size=Pt(13.5)), lh=1.02)])
        text(s, x + pad, top + Pt(118), cw - pad * 2, h - Pt(130),
             [para(run(desc, font=SANS, color=GRAY, size=Pt(11)), lh=1.16)])
        if i < n - 1:
            text(s, x + cw, top + h / 2 - Pt(12), gap, Pt(24),
                 [para(run("\uf061", font=ICON, color=ORANGE, size=Pt(11)),
                       align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
        x += cw + gap
    by = top + h + Pt(30)
    rect(s, MARGIN, by, CONTENT_W, Pt(58), fill=DEEP, radius=0.1)
    icon_token(s, MARGIN + Pt(24), by + Pt(14), Pt(30), FA["check"], fg=AMBER,
               bg=None)
    text(s, MARGIN + Pt(66), by, CONTENT_W - Pt(100), Pt(58),
         [para([("Takeaway.  ", {"font": DISPLAY, "color": WHITE,
                                 "size": Pt(13)}),
                ("Uncompromised privacy-cap enforcement at a fraction of the "
                 "memory footprint.", {"font": SANS, "color":
                                       RGBColor(0xD7, 0xDB, 0xE2),
                                       "size": Pt(12.5)})])],
         anchor=MSO_ANCHOR.MIDDLE)
    footer(s, page)
    return s


def s_distortion(prs, page):
    s = s_blank(prs)
    bg(s)
    top = header(
        s, "Distortion visibility",
        "We quantify the true cost of meeting privacy caps.",
        "Every percentage point of distortion is measured \u2014 not assumed.")
    cx, cy = MARGIN, top + Pt(2)
    cw, ch = Pt(486), Pt(296)
    rect(s, cx, cy, cw, ch, fill=PANEL, line=HAIR, line_w=Pt(1), radius=0.04)
    # plot area
    pad_l, pad_t, pad_b, pad_r = Pt(40), Pt(54), Pt(46), Pt(28)
    ax0 = cx + pad_l
    ay0 = cy + ch - pad_b
    pw = cw - pad_l - pad_r
    phh = ch - pad_t - pad_b
    # gridlines
    for g in range(1, 4):
        gy = cy + pad_t + phh * g / 4
        line_seg(s, ax0, gy, ax0 + pw, gy, HAIR2, Pt(0.75))
    line_seg(s, ax0, cy + pad_t, ax0, ay0, HAIR, Pt(1))
    line_seg(s, ax0, ay0, ax0 + pw, ay0, HAIR, Pt(1))

    raw = [0.30, 0.55, 0.42, 0.66, 0.50, 0.72, 0.58, 0.40]
    bal = [0.16, 0.34, 0.28, 0.40, 0.32, 0.44, 0.36, 0.24]

    def pt(i, v, series):
        xx = ax0 + pw * i / (len(series) - 1)
        yy = ay0 - phh * v
        return (int(xx), int(yy))

    raw_pts = [pt(i, v, raw) for i, v in enumerate(raw)]
    bal_pts = [pt(i, v, bal) for i, v in enumerate(bal)]
    # shaded distortion band
    band = raw_pts + list(reversed(bal_pts))
    area = free_poly(s, band, fill=ORANGE_SOFT)
    _alpha(area, 70)
    # delta connectors at sample points
    for rp, bp in zip(raw_pts, bal_pts):
        line_seg(s, rp[0], rp[1], bp[0], bp[1], ORANGE, Pt(0.75), dash="dash")
    # strokes
    for a, b in zip(raw_pts, raw_pts[1:]):
        line_seg(s, a[0], a[1], b[0], b[1], GRAY, Pt(1.5), dash="dash")
    for a, b in zip(bal_pts, bal_pts[1:]):
        line_seg(s, a[0], a[1], b[0], b[1], ORANGE, Pt(2.4))
    for p_ in bal_pts:
        oval(s, p_[0] - Pt(3), p_[1] - Pt(3), Pt(6), Pt(6), ORANGE)
    # labels
    text(s, cx + Pt(20), cy + Pt(16), Pt(200), Pt(18),
         [para([("\u2509 ", {"font": SANS_SB, "color": GRAY, "size": Pt(11)}),
                ("Raw market share", {"font": SANS_SB, "color": GRAY,
                                      "size": Pt(10.5)})])])
    text(s, cx + Pt(20), cy + Pt(32), Pt(220), Pt(18),
         [para([("\u2501 ", {"font": SANS_SB, "color": ORANGE, "size": Pt(11)}),
                ("Balanced (compliant) share", {"font": SANS_SB,
                 "color": ORANGE_DK, "size": Pt(10.5)})])])
    text(s, ax0 + pw / 2 - Pt(60), ay0 + Pt(12), Pt(160), Pt(16),
         [para(run("\u0394 distortion measured per category", font=SANS,
                   color=MUTE, size=Pt(9)), align=PP_ALIGN.CENTER)])

    rx = cx + cw + Pt(28)
    rw = SLIDE_W - MARGIN - rx
    blocks = [
        (FA["chart"], "Distortion Analysis", "Computes the exact percentage-"
         "point delta between raw market truth and compliant reporting."),
        (FA["table"], "Rank-Change Tracking", "A dedicated sheet records every "
         "rank shift before and after weighting."),
        (FA["scale"], "Preset Comparison", "--compare-presets builds a matrix "
         "of mean distortion to find the optimal configuration."),
    ]
    by = top + Pt(2)
    bh = Pt(74)
    for glyph, t, b in blocks:
        rect(s, rx, by, rw, bh, fill=WHITE, line=HAIR, line_w=Pt(1),
             radius=0.06, shadow=True)
        icon_token(s, rx + Pt(16), by + (bh - Pt(34)) / 2, Pt(34), glyph,
                   fg=ORANGE)
        text(s, rx + Pt(62), by + Pt(12), rw - Pt(76), Pt(20),
             [para(run(t, font=DISPLAY, color=INK, size=Pt(13)))])
        text(s, rx + Pt(62), by + Pt(33), rw - Pt(76), bh - Pt(40),
             [para(run(b, font=SANS, color=GRAY, size=Pt(10.5)), lh=1.12)])
        by += bh + Pt(11)
    rect(s, rx, by, rw, Pt(56), fill=ORANGE_SOFT, line=ORANGE, line_w=Pt(1),
         radius=0.08)
    text(s, rx + Pt(18), by, rw - Pt(34), Pt(56),
         [para([("Insight.  ", {"font": DISPLAY, "color": ORANGE_DK,
                                "size": Pt(11.5)}),
                ("We don\u2019t just enforce rules \u2014 we quantify their "
                 "mathematical impact.", {"font": SANS, "color": INK2,
                                          "size": Pt(11)})], lh=1.14)],
         anchor=MSO_ANCHOR.MIDDLE)
    footer(s, page)
    return s


def s_policy(prs, page):
    s = s_blank(prs)
    bg(s)
    top = header(
        s, "Policy enforcement",
        "The system actively guards against unauthorized data combinations.")
    rows = [
        (FA["check"], "Fraud / Chargebacks", "Must set privacy_basis: "
         "clearing_spend", "ENFORCED", DEEP, False),
        (FA["warn"], "Digital Wallets", "Requires explicit Privacy-team "
         "review", "MANUAL APPROVAL", ORANGE_DK, False),
        (FA["warn"], "Dual-Entity-Axis", "Requires explicit Privacy-team "
         "review", "MANUAL APPROVAL", ORANGE_DK, False),
        (FA["ban"], "Top-Merchant Lists", "Strictly prohibited",
         "HARD BLOCKED", INK, True),
    ]
    headers = ["", "Data scenario", "Policy condition", "Status"]
    widths = [Pt(48), Pt(250), Pt(430), Pt(168)]
    x = MARGIN
    y = top + Pt(4)
    head_h = Pt(34)
    cxx = x
    for htext, cwid in zip(headers, widths):
        if htext:
            text(s, cxx + Pt(4), y, cwid - Pt(8), head_h,
                 [para(run(htext.upper(), font=SANS_SB, color=GRAY,
                           size=Pt(10), track=0.8))], anchor=MSO_ANCHOR.MIDDLE)
        cxx += cwid
    line_seg(s, x, y + head_h, x + CONTENT_W, y + head_h, INK, Pt(1.6))
    ry = y + head_h
    row_h = Pt(60)
    for glyph, scenario, cond, badge, bcol, solid in rows:
        icon_token(s, x + Pt(2), ry + (row_h - Pt(32)) / 2, Pt(32), glyph,
                   fg=bcol if not solid else INK,
                   bg=ORANGE_SOFT if bcol == ORANGE_DK else DEEP_SOFT)
        text(s, x + widths[0] + Pt(4), ry, widths[1] - Pt(8), row_h,
             [para(run(scenario, font=DISPLAY, color=INK, size=Pt(13)))],
             anchor=MSO_ANCHOR.MIDDLE)
        text(s, x + widths[0] + widths[1] + Pt(4), ry, widths[2] - Pt(8),
             row_h,
             [para(run(cond, font=MONO, color=INK2, size=Pt(11)))],
             anchor=MSO_ANCHOR.MIDDLE)
        bx = x + widths[0] + widths[1] + widths[2] + Pt(4)
        chip(s, bx, ry + (row_h - Pt(24)) / 2, badge, bcol, bcol,
             solid=solid, mono=True)
        line_seg(s, x, ry + row_h, x + CONTENT_W, ry + row_h, HAIR2, Pt(1))
        ry += row_h
    footer(s, page)
    return s


def s_outputs(prs, page):
    s = s_blank(prs)
    bg(s)
    top = header(
        s, "Deliverables",
        "Actionable intelligence, packaged for decisions and for audit.")
    cards = [
        (FA["excel"], "The Excel Report", ".xlsx", [
            "Executive summary & run metadata.",
            "Category comparisons (Target vs. Best-in-Class).",
            "Rank changes & weight methods.",
        ]),
        (FA["csv"], "The Balanced Export", ".csv", [
            "BI-ready for Tableau & Power BI.",
            "Enriched metrics (Distortion_PP, Raw_Share).",
            "Schema-validated for parity with Excel.",
        ]),
        (FA["archive"], "The Audit Package", ".zip", [
            "Immutable snapshot of inputs, config & logs.",
            "Ready for immediate regulatory submission.",
        ]),
    ]
    h = Pt(300)
    gap = Pt(24)
    n = len(cards)
    cw = Emu(int((CONTENT_W - gap * (n - 1)) / n))
    x = MARGIN
    for glyph, title, ext, items in cards:
        rect(s, x, top, cw, h, fill=WHITE, line=HAIR, line_w=Pt(1),
             radius=0.05, shadow=True)
        pad = Pt(22)
        icon_token(s, x + pad, top + Pt(24), Pt(44), glyph, fg=ORANGE)
        text(s, x + cw - Pt(110), top + Pt(30), Pt(90) - Pt(0), Pt(30),
             [para(run(ext, font=MONO, color=ORANGE_DK, size=Pt(13)),
                   align=PP_ALIGN.RIGHT)])
        text(s, x + pad, top + Pt(82), cw - pad * 2, Pt(34),
             [para(run(title, font=DISPLAY, color=INK, size=Pt(16.5)))])
        line_seg(s, x + pad, top + Pt(120), x + cw - pad, top + Pt(120),
                 HAIR2, Pt(1))
        text(s, x + pad, top + Pt(134), cw - pad * 2, h - Pt(150),
             bullets([run(i, color=INK2) for i in items], size=Pt(12),
                     sa=Pt(10)))
        x += cw + gap
    footer(s, page)
    return s


def s_zerotrust(prs, page):
    s = s_blank(prs)
    bg(s)
    top = header(
        s, "Verified execution",
        "A zero-trust release pipeline guarantees mathematical integrity.")
    nodes = [(FA["branch"], "Code Commit"), (FA["gears"], "GitHub Actions"),
             (FA["clipboard"], "Gate Tests"), (FA["check"], "Validated Release")]
    nw = Pt(190)
    nh = Pt(64)
    gap = Pt(40)
    total = nw * len(nodes) + gap * (len(nodes) - 1)
    x = MARGIN + (CONTENT_W - total) / 2
    y = top + Pt(6)
    for i, (glyph, label) in enumerate(nodes):
        last = i == len(nodes) - 1
        rect(s, x, y, nw, nh, fill=WHITE if not last else ORANGE_SOFT,
             line=ORANGE if last else HAIR, line_w=Pt(1.4) if last else Pt(1),
             radius=0.14, shadow=True)
        icon_token(s, x + Pt(14), y + (nh - Pt(34)) / 2, Pt(34), glyph,
                   fg=ORANGE)
        text(s, x + Pt(56), y, nw - Pt(64), nh,
             [para(run(label, font=DISPLAY, color=INK, size=Pt(13)))],
             anchor=MSO_ANCHOR.MIDDLE)
        if not last:
            text(s, x + nw, y, gap, nh,
                 [para(run("\uf061", font=ICON, color=ORANGE, size=Pt(12)),
                       align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
        x += nw + gap
    text(s, MARGIN, y + nh + Pt(8), CONTENT_W, Pt(16),
         [para([("\uf01e  ", {"font": ICON, "color": MUTE, "size": Pt(8)}),
                ("continuous validation loop", {"font": SANS, "color": MUTE,
                 "size": Pt(9.5), "italic": True})], align=PP_ALIGN.CENTER)])

    details = [
        (FA["clipboard"], "Gate Test Suite",
         "17+ representative scenarios run on every commit."),
        (FA["scale"], "Mathematical Parity",
         "CSV exports are cross-checked against compliant Excel outputs to a "
         "0.01% tolerance."),
        (FA["lock"], "Immutable Offline Deploys",
         "Checksum-verified bundles (offline_packages/) for air-gapped "
         "environments."),
    ]
    ty = y + nh + Pt(40)
    th = Pt(132)
    gap2 = Pt(24)
    cw = Emu(int((CONTENT_W - gap2 * 2) / 3))
    xx = MARGIN
    for glyph, t, b in details:
        rect(s, xx, ty, cw, th, fill=PANEL, line=HAIR, line_w=Pt(1),
             radius=0.06)
        icon_token(s, xx + Pt(18), ty + Pt(18), Pt(34), glyph, fg=ORANGE)
        text(s, xx + Pt(18), ty + Pt(60), cw - Pt(36), Pt(24),
             [para(run(t, font=DISPLAY, color=INK, size=Pt(13)))])
        text(s, xx + Pt(18), ty + Pt(84), cw - Pt(36), th - Pt(94),
             [para(run(b, font=SANS, color=GRAY, size=Pt(11)), lh=1.16)])
        xx += cw + gap2
    footer(s, page)
    return s


def s_value(prs, page):
    s = s_blank(prs)
    bg(s)
    top = header(
        s, "Value realization",
        "Autobench turns regulatory constraint into analytical advantage.")
    cols = [
        (FA["shield"], "Absolute Compliance", [
            "100% adherence to Control 3.2 privacy caps.",
            "Cryptographically verifiable audit trails.",
        ]),
        (FA["gauge"], "Mathematical Superiority", [
            "Industry-leading LP solvers minimize data distortion.",
            "Unlocks insight in sparse data that legacy tools abandon.",
        ]),
        (FA["bolt"], "Operational Agility", [
            "TUI and lean mode democratize access at massive scale.",
            "No infrastructure bloat.",
        ]),
    ]
    h = Pt(214)
    gap = Pt(24)
    n = len(cols)
    cw = Emu(int((CONTENT_W - gap * (n - 1)) / n))
    x = MARGIN
    for glyph, title, items in cols:
        rect(s, x, top, cw, h, fill=PANEL, line=HAIR, line_w=Pt(1),
             radius=0.05)
        icon_token(s, x + Pt(20), top + Pt(22), Pt(42), glyph, fg=ORANGE)
        text(s, x + Pt(20), top + Pt(78), cw - Pt(40), Pt(48),
             [para(run(title, font=DISPLAY, color=INK, size=Pt(16)), lh=1.02)])
        text(s, x + Pt(20), top + Pt(126), cw - Pt(40), h - Pt(140),
             bullets([run(i, color=INK2) for i in items], size=Pt(12),
                     sa=Pt(9)))
        x += cw + gap
    by = top + h + Pt(30)
    rect(s, MARGIN, by, CONTENT_W, Pt(68), fill=ORANGE, radius=0.1)
    text(s, MARGIN, by, CONTENT_W, Pt(68),
         [para(run("Compliant data is no longer a compromise. "
                   "It is an exact science.", font=DISPLAY, color=WHITE,
                   size=Pt(18)), align=PP_ALIGN.CENTER)],
         anchor=MSO_ANCHOR.MIDDLE)
    footer(s, page)
    return s


def s_closing(prs):
    s = s_blank(prs)
    bg(s, DARK)
    d = Pt(120)
    sep = d * 0.52
    total = sep + d
    cxp = (SLIDE_W - total) / 2
    cyp = (SLIDE_H - d) / 2 - Pt(26)
    oval(s, cxp, cyp, d, d, RED)
    oval(s, cxp + sep, cyp, d, d, AMBER)
    oval(s, cxp + sep, cyp + (d - d * 0.855) / 2, d - sep, d * 0.855, OVERLAP)
    text(s, MARGIN, cyp + d + Pt(20), CONTENT_W, Pt(40),
         [para(run("mastercard", font=DISPLAY, color=WHITE, size=Pt(24),
                   track=0.5), align=PP_ALIGN.CENTER)])
    text(s, MARGIN, cyp + d + Pt(64), CONTENT_W, Pt(24),
         [para(run("AUTOBENCH v3.0  \u2022  PRECISION  \u2022  COMPLIANCE  "
                   "\u2022  CONTROL", font=SANS_SB,
                   color=RGBColor(0x9A, 0x9E, 0xA6), size=Pt(11), track=2.2),
               align=PP_ALIGN.CENTER)])
    return s


# --------------------------------------------------------------------------- #
# Assembly
# --------------------------------------------------------------------------- #
def build(output: Path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    s_title(prs)
    s_paradox(prs, 2)
    s_architecture(prs, 3)
    s_solver(prs, 4)
    s_caps(prs, 5)
    s_presets(prs, 6)
    s_tui(prs, 7)
    s_resource(prs, 8)
    s_distortion(prs, 9)
    s_policy(prs, 10)
    s_outputs(prs, 11)
    s_zerotrust(prs, 12)
    s_value(prs, 13)
    s_closing(prs)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    n = embed_fonts(output, FONT_DIR, EMBED_MAP, theme_major=DISPLAY,
                    theme_minor=SANS)
    print(f"Saved {output} ({len(prs.slides._sldIdLst)} slides, "
          f"{n} embedded font parts)")


def main():
    out = Path(sys.argv[1]) if len(sys.argv) > 1 else Path(
        "docs/Autobench_v3.0_Precision_Compliance.pptx")
    build(out)


if __name__ == "__main__":
    main()
