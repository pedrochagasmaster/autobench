"""Add presenter speaker notes to the Autobench v3.0 deck.

Notes are sourced verbatim from
``docs/Autobench_v3.0_Precision_Compliance_PRESENTATION_SCRIPT.md`` and mapped to
the live PPTX slide order (15 slides; four detail slides sit between the eleven
slides that appear in the presented PDF).

Usage::

    py scripts/add_pptx_notes.py [input.pptx] [output.pptx]

If ``output`` is omitted the input file is updated in place.
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation

# --------------------------------------------------------------------------- #
# Notes keyed by 1-based slide index in the live PPTX.
# --------------------------------------------------------------------------- #
NOTES: dict[int, str] = {
    # 1 - Title -------------------------------------------------------------- #
    1: """OPENING (deliver before advancing, or as the first lines on this slide):
"Today I will walk through Autobench v3.0, our precision compliance engine for peer benchmarking. The simple version is this: Autobench lets us compare an entity against a peer group while enforcing privacy rules mathematically. It protects peer confidentiality, preserves as much market signal as possible, and creates the evidence trail needed to defend the result."

SPEAKER SCRIPT:
"Autobench exists because peer benchmarking has become too important to run as a manual balancing exercise.

When we benchmark an issuer, acquirer, merchant, or market segment, we need two things at the same time. The result has to be analytically useful, and it has to protect the confidentiality of every peer in the comparison.

Those requirements can pull against each other. If we enforce privacy too bluntly, we lose market truth. If we preserve market truth too casually, we risk violating concentration caps.

Autobench is the engine built for that exact tension. It takes raw peer data, applies Control 3 privacy constraints, solves the balancing problem mathematically, and produces outputs that analysts can explain and auditors can review.

So the promise of this deck is the subtitle: next-generation benchmarking, minimal distortion, and absolute control."

EMPHASIZE: "engine" and "mathematically" - this frames Autobench as governed infrastructure, not a convenience script.
TRANSITION: "The reason this matters starts with the compliance paradox.\"""",
    # 2 - Compliance Paradox ------------------------------------------------- #
    2: """SPEAKER SCRIPT:
"This is the core problem.

Control 3 constraints exist to prevent a benchmark from exposing or over-weighting any individual peer. We need concentration caps. We need minimum participant thresholds. And we need confidence that the published benchmark is not only useful, but compliant.

At the same time, the business needs market reality. The benchmark has to preserve positioning, rank integrity, and performance differences that matter for commercial decisions.

That creates the paradox. The privacy rules protect the benchmark, but applying them manually can distort the benchmark.

In a manual workflow, an analyst adjusts peer weights, checks the caps, adjusts again, checks again, and then has to explain why the final result is still representative. That process is slow, inconsistent, and hard to audit.

Autobench changes the operating model. It treats compliance as a constrained optimization problem. The engine searches for weights that satisfy the privacy caps while keeping the balanced result as close as possible to the raw market data.

That is the shift: compliance is not a final inspection step. It is embedded in the calculation itself."

EMPHASIZE: the phrase "not a final inspection step" - the main executive takeaway.
TRANSITION: "The architecture is built around that idea: ingest, enforce, and publish evidence.\"""",
    # 3 - Architecture ------------------------------------------------------- #
    3: """SPEAKER SCRIPT:
"The architecture is deliberately simple at the top level.

Phase one is ingestion. The tool can be used through automated pipelines, the CLI, or the analyst workbench. The goal is to bring the data into a normalized, validated structure before any benchmark is calculated.

Phase two is the core engine. This is where Autobench applies schema validation, Control 3.2 privacy validation, and global linear-programming optimization. The important point is that the solver is not calculating a simple peer average. It is solving for compliant peer weights across the categories in scope.

Phase three is output. Autobench produces Excel reports, diagnostic sheets, balanced CSV exports, and audit packages. That turns the analysis from a result into a defensible artifact.

This gives us a clean chain of custody. We know what data went in, what privacy rules were applied, how the weights were chosen, and what evidence was produced for review."

EMPHASIZE: "chain of custody" - it connects architecture to auditability.
TRANSITION: "The next question is how analysts interact with that engine in practice.\"""",
    # 4 - Algorithmic Resilience (detail / not in PDF flow) ------------------ #
    4: """BACKUP / DETAIL SLIDE - not part of the 11-slide PDF flow the presentation script covers.

Skip in the standard 10-12 minute talk. Use on demand if the audience drills into HOW the solver guarantees an answer: the three-tier cascade of Global LP optimization, then subset search with per-dimension LP, then a Bayesian / heuristic (L-BFGS-B) fallback for sparse edge cases.

If you do present it, the one-liner is: "Strict optimization is tried first; structured fallbacks guarantee a usable, still-compliant result even when the data is hard.\"""",
    # 5 - Control 3.2 caps table (detail / not in PDF flow) ------------------ #
    5: """BACKUP / DETAIL SLIDE - not part of the 11-slide PDF flow the presentation script covers.

Skip in the standard talk. Use on demand if someone asks exactly how the privacy rule is chosen: the engine auto-selects the cap rule from the peer count (5/25, 6/30, 7/35, 10/40, and 4/35 for merchant benchmarking), and higher tiers add participation requirements. No manual configuration.

Reinforces the Q&A answer "Does Autobench ever bypass a privacy cap?" - it does not; caps are constraints, and infeasibility is surfaced rather than bypassed.""",
    # 6 - Analyst Experience ------------------------------------------------- #
    6: """SPEAKER SCRIPT:
"The analyst experience matters because compliance only scales if the workflow is repeatable.

The TUI gives analysts a guided path through the run. First, it is validation-first. Before the heavy computation starts, it checks for nulls, schema errors, and peer-constraint issues. That catches bad inputs before they become bad outputs.

Second, configuration is controlled. Analysts select the file, entity, metrics, dimensions, presets, and output mode in a structured interface rather than rebuilding a command from memory.

Third, diagnostics are visible while the run is happening. The analysis executes in the background, the interface remains responsive, and logs are surfaced in the tool.

Fourth, advanced overrides are still available for expert users. The difference is that those overrides live inside a governed workflow.

The business benefit is consistency. The same analysis intent produces the same behavior, even when different analysts run it."

PRESENTER NOTE: if someone asks how expert overrides are exposed, the TUI's advanced panel toggles on with Ctrl+A. The point is that expert controls live inside the guided workflow, not beside it.
EMPHASIZE: the TUI is not only usability; it is standardization.
TRANSITION: "Once the benchmark is produced, the next question is how much the privacy balancing moved the market signal.\"""",
    # 7 - Resource Management (detail / not in PDF flow) --------------------- #
    7: """BACKUP / DETAIL SLIDE - not part of the 11-slide PDF flow the presentation script covers.

Skip in the standard talk. Use on demand if the audience asks how Autobench runs in memory-limited or large-data environments: it estimates and projects memory load, streams heavy CSVs in chunks, pre-aggregates duplicate rows, and offers a lean execution mode that focuses compute on the global optimization.

Ties to the "Operational Agility" pillar on the Value slide ("lean execution keeps it practical in memory-limited environments").""",
    # 8 - Distortion Visibility ---------------------------------------------- #
    8: """SPEAKER SCRIPT:
"This is one of the most important parts of the tool.

Autobench does not just say, 'the benchmark is compliant.' It also tells us what compliance cost analytically.

For each category, the tool compares raw market share with the balanced, compliant share. The difference is reported as distortion in percentage points. That gives analysts a direct way to answer: how much did the privacy constraint move the result?

The tool also tracks rank changes. If weighting changes the ordering of peers or categories, that is recorded. That matters because rank movement is often where stakeholders first challenge the result.

Preset comparison adds another layer. Analysts can compare configurations and see which one produces the lowest mean distortion while still respecting the rules.

The takeaway is that Autobench makes the trade-off visible. It does not hide the cost of compliance inside a black box. It exposes the cost so the team can make a better decision and defend that decision later."

EMPHASIZE: "what compliance cost analytically" - the sharpest framing for this slide.
TRANSITION: "Mathematical compliance is necessary, but it is not the whole policy picture. Some cases require governance controls.\"""",
    # 9 - Policy Enforcement ------------------------------------------------- #
    9: """SPEAKER SCRIPT:
"This slide is about policy controls beyond the numeric cap calculation.

Some requirements can be enforced directly by the system. Fraud and chargeback analysis is an example. The run must declare the correct privacy basis: clearing spend. If that condition is not present, the workflow blocks the run.

Some requirements need explicit Privacy review. Digital wallet analysis and dual-entity-axis analysis fall into that category. The tool cannot replace the governance decision, but it can require evidence that the decision happened.

And some outputs are simply not allowed. Top-merchant lists are hard-blocked because they are outside the permitted policy posture.

That separation is important. Autobench does not pretend every compliance question can be reduced to a formula. It separates enforceable controls, manual approval gates, and prohibited outputs.

That is what makes the workflow credible. The tool automates what should be automated, and it forces explicit review where human governance is required."

EMPHASIZE: "automates what should be automated" - this avoids overclaiming.
TRANSITION: "Once the run clears those controls, the output has to serve multiple audiences.\"""",
    # 10 - Deliverables ------------------------------------------------------ #
    10: """SPEAKER SCRIPT:
"Autobench produces three deliverable types.

The Excel report is the primary decision artifact. It gives stakeholders the summary, run metadata, category comparisons, target versus best-in-class views, rank changes, and weight methods.

The balanced CSV is the data product. It is designed for Tableau, Power BI, and other downstream workflows. It can carry raw share, balanced share, distortion, and other calculated metrics. Just as important, it is validated for parity with the Excel report.

The audit package is the evidence layer. It captures inputs, configuration, and logs so the run can be reproduced and reviewed.

The point is that Autobench does not stop at calculation. It packages the result for the people who need to use it: executives who need a conclusion, analysts who need data, and reviewers who need evidence."

EMPHASIZE: "decision artifact, data product, evidence layer."
TRANSITION: "Those deliverables are what convert the engine into business value.\"""",
    # 11 - Verified Execution (detail / not in PDF flow) --------------------- #
    11: """BACKUP / DETAIL SLIDE - not part of the 11-slide PDF flow the presentation script covers.

Skip in the standard talk. Use on demand if the audience asks how the tool itself is trusted: a gate-test suite of representative scenarios runs on every commit, CSV exports are cross-checked against the compliant Excel outputs, and offline deploys are checksum-verified.

CAUTION: this slide leans on strong wording. Pair it with the honest framing from the Value slide and Q&A - verification today is reproducibility and parity-checking, not cryptographic signing.""",
    # 12 - Value Realization ------------------------------------------------- #
    12: """SPEAKER SCRIPT:
"The value realization comes from three places.

First, absolute compliance. The engine is built around Control 3.2 privacy caps, and the output trail is designed for review. That lowers the risk of inconsistent manual application.

Second, mathematical superiority. The solver minimizes distortion instead of relying on manual balancing judgment. That lets us keep more of the market signal while still enforcing the privacy rules.

Third, operational agility. The TUI makes the workflow accessible. Lean execution keeps it practical in memory-limited environments. And the deliverables are ready for both stakeholder review and downstream analysis.

So the larger point is this: compliant data is no longer a compromise. If we can measure the constraint, optimize around it, and document the result, compliance becomes part of analytical quality."

PRESENTER NOTE: the slide reads "cryptographically verifiable audit trails." The script deliberately says the output trail is "designed for review" because today the audit package is verified through completeness and reproducibility, not a cryptographic signature. If anyone presses on the slide wording, use the audit-trail answer in Q&A instead of defending a signing claim.
EMPHASIZE: the final sentence - it is the strategic claim of the presentation.
TRANSITION: "The next slide quantifies the operational side of that value.\"""",
    # 13 - Efficiency & Scalability ------------------------------------------ #
    13: """SPEAKER SCRIPT:
"This slide shows the operating impact.

The extraction step still exists. We still need the right source data and the right peer definitions. The major gain is in balancing, which was previously the most manual and senior-dependent part of the process.

Before Autobench, balancing one peer could take roughly five to seven days. With the tool, that becomes roughly one to two days. On average, that is about a 75% reduction in balancing time per peer.

The effect compounds when we move from one peer to a market-wide study. A five-peer project that previously required about 30 days of balancing effort can move toward about 7.5 days. A six-peer project can move from roughly 36 days to about 9 days.

That changes the operating model. Large benchmark studies no longer need to be treated as bespoke spreadsheet exercises. They can be executed as repeatable analytical runs, with senior time redirected toward interpretation, client strategy, and higher-value advisory work."

NUMBERS TO KNOW COLD: ~5-7 days -> ~1-2 days per peer (avg 6 -> 1.5); ~75% average reduction; 5 peers ~30 -> ~7.5 days; 6 peers ~36 -> ~9 days. These are balancing-phase gains; extraction time is unchanged.
EMPHASIZE: this is not just cycle-time reduction; it is a staffing and scale change.
TRANSITION: "The final content slide shows what this looks like in a field use case.\"""",
    # 14 - A&F Dashboard ----------------------------------------------------- #
    14: """A&F = Authorization & Fraud (auth rates and fraud rates are the dashboard's core metric families). Confirm the house definition before presenting.

SPEAKER SCRIPT:
"The A&F dashboard is where the story becomes concrete.

The first foundation is peer readiness. The dashboard uses pre-defined peer groups for issuers and acquirers, with 14 months of history from January 2025 through February 2026. Brazil is live, and the broader LAC rollout is underway.

The second foundation is metric depth. The dashboard can support raw and clean authorization rates, fraud rates, declined-reason share, count and amount views, credit and debit splits, card-present and card-not-present cuts, domestic and cross-border cuts, tokenization, 3DS, MCC, ticket size, and more.

The third foundation is ACS impact. Autobench reduces effort on the hardest parts of the workflow: extraction, peer grouping, compliant balancing, and repeatable output production. That frees senior people to focus on diagnosis, recommendations, and client value.

This also opens the door to a stronger commercial model: A&F Report as a Service. Instead of rebuilding each benchmark manually, we can scale a governed benchmarking engine across markets and peer groups.

That is the broader significance of Autobench. It is not only improving one analysis process. It creates reusable infrastructure for privacy-compliant performance analytics."

CLOSING SCRIPT:
"To close: Autobench solves the core tension in peer benchmarking. It protects privacy, preserves market signal, quantifies distortion, creates audit evidence, and reduces the manual effort required to scale. That is why it is best thought of as a precision compliance engine, not just a reporting tool."

TRANSITION TO Q&A: "I will pause there. I am happy to go deeper on the optimization method, the Control 3 policy gates, or the A&F dashboard rollout.\"""",
    # 15 - Blank final page -------------------------------------------------- #
    15: """BLANK FINAL PAGE - no spoken content is required.

If the blank page appears, use it as a neutral Q&A screen:
"I will pause there and take questions.\"""",
}

# Keyword each slide's on-slide text must contain, as a guard against the slide
# order drifting from the mapping above.
EXPECTED = {
    1: "PRECISION COMPLIANCE",
    2: "COMPLIANCE PARADOX",
    3: "ARCHITECTURE",
    4: "ALGORITHMIC RESILIENCE",
    5: "CONTROL 3.2",
    6: "ANALYST EXPERIENCE",
    7: "RESOURCE MANAGEMENT",
    8: "DISTORTION VISIBILITY",
    9: "POLICY ENFORCEMENT",
    10: "DELIVERABLES",
    11: "VERIFIED EXECUTION",
    12: "VALUE REALIZATION",
    13: "SCALABILITY",
    14: "USE CASE",
    15: None,
}


def _slide_text(slide) -> str:
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            parts.append(shape.text_frame.text)
    return " ".join(parts).upper()


def main() -> int:
    in_path = Path(sys.argv[1]) if len(sys.argv) > 1 else Path(
        "docs/Autobench_v3.0_Precision_Compliance.pptx")
    out_path = Path(sys.argv[2]) if len(sys.argv) > 2 else in_path

    prs = Presentation(str(in_path))
    slides = list(prs.slides)
    if len(slides) != 15:
        print(f"WARNING: expected 15 slides, found {len(slides)}")

    applied = 0
    for idx, slide in enumerate(slides, start=1):
        keyword = EXPECTED.get(idx)
        if keyword and keyword not in _slide_text(slide):
            raise SystemExit(
                f"Slide {idx} guard failed: expected to contain {keyword!r}; "
                f"got: {_slide_text(slide)[:120]!r}")
        note = NOTES.get(idx)
        if not note:
            continue
        slide.notes_slide.notes_text_frame.text = note
        applied += 1

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"Applied notes to {applied} slides -> {out_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
